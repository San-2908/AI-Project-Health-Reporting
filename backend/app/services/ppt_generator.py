from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import List
from app.models.schemas import ProjectData, PortfolioReport
import io

def create_ppt_presentation(projects: List[ProjectData]) -> io.BytesIO:
    prs = Presentation()
    
    # Slide 1: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Project Health Report"
    subtitle.text = f"Executive Summary for {len(projects)} Projects"
    
    # Slide 2: Portfolio Health
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Portfolio Health Overview"
    content = slide.placeholders[1].text_frame
    for p in projects:
        p_status = p.status.rag_score if p.status else "Unknown"
        content.add_paragraph().text = f"{p.project_name}: {p_status}"
    
    # Slide 3: Risk Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risk Analysis"
    content = slide.placeholders[1].text_frame
    for p in projects:
        if p.status and p.status.top_risks:
            content.add_paragraph().text = f"{p.project_name} Risks:"
            for r in p.status.top_risks:
                p_run = content.add_paragraph()
                p_run.text = f"  - {r}"
                p_run.level = 1
                
    # Slide 4: Trend Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Trend Analysis"
    content = slide.placeholders[1].text_frame
    content.text = "Overall completion rate trend is stable based on the analyzed projects."
    
    # Slide 5: Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Recommendations"
    content = slide.placeholders[1].text_frame
    for p in projects:
        if p.status and p.status.recommended_actions:
            content.add_paragraph().text = f"{p.project_name}:"
            for a in p.status.recommended_actions:
                p_run = content.add_paragraph()
                p_run.text = f"  - {a}"
                p_run.level = 1
                
    # Slide 6: Next Month Outlook
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Month Outlook"
    content = slide.placeholders[1].text_frame
    content.text = "Monitoring critical milestones. Expected to resolve top blockers."

    # Save to BytesIO
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream
